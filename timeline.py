import streamlit as st
import datetime
import json
import base64
import importlib.util
import subprocess
import sys
from io import BytesIO
from uuid import uuid4
import pandas as pd
try:
    from spellchecker import SpellChecker
    SPELLCHECKER_AVAILABLE = True
except ImportError:
    SPELLCHECKER_AVAILABLE = False
    SpellChecker = None
import streamlit.components.v1 as components
import os
import zipfile
from pathlib import Path


def rerun_app():
    rerun_fn = getattr(st, "rerun", None) or getattr(st, "experimental_rerun", None)
    if rerun_fn is None:
        raise AttributeError("Streamlit rerun function not available")
    rerun_fn()

DEFAULT_TITLE = "Interactive Event Timeline with Lens Magnifier"
VALID_VIEWS = {"Timeline", "Domino"}


def normalize_view_choice(value: str) -> str:
    """Map stored view names to the currently supported set."""
    if not value:
        return "Timeline"
    if value == "Orbit":
        return "Domino"
    return value if value in VALID_VIEWS else "Timeline"

@st.cache_resource
def get_spell_checker():
    """Initialize and cache a spell checker instance."""
    if not SPELLCHECKER_AVAILABLE:
        return None
    spell = SpellChecker()
    # Add common event-related words that might not be in the dictionary
    custom_words = {
        'streamlit', 'timeline', 'event', 'magnifier', 'lens', 'orbit', 
        'visualization', 'interactive', 'dataset', 'webinar', 'workshop',
        'conference', 'meeting', 'presentation', 'deadline', 'milestone',
        'launch', 'release', 'update', 'announcement', 'celebration',
        'birthday', 'anniversary', 'holiday', 'vacation', 'travel',
        'project', 'sprint', 'retrospective', 'planning', 'review'
    }
    spell.word_frequency.load_words(custom_words)
    return spell

def check_spelling_and_suggest(text, spell_checker):
    """Check spelling of text and return suggestions for misspelled words."""
    if not spell_checker or not SPELLCHECKER_AVAILABLE:
        return [], []
    
    if not text or not text.strip():
        return [], []
    
    words = text.split()
    misspelled = spell_checker.unknown([word.lower().strip('.,!?;:"\'()[]{}') for word in words])
    suggestions = {}
    
    for word in misspelled:
        candidates = spell_checker.candidates(word)
        if candidates:
            suggestions[word] = list(candidates)[:3]  # Get top 3 suggestions
    
    return misspelled, suggestions

def display_spell_suggestions(text, misspelled, suggestions):
    """Display spelling suggestions in a user-friendly format."""
    if not misspelled:
        return st.success("✅ No spelling errors found!")
    
    st.warning(f"⚠️ Found {len(misspelled)} potential spelling error(s):")
    
    for word in misspelled:
        if word in suggestions:
            suggestion_text = ", ".join(suggestions[word])
            st.info(f"• **{word}** - Did you mean: {suggestion_text}?")
        else:
            st.info(f"• **{word}** - No suggestions available")

def save_to_localstorage():
    """Generate JavaScript to save current session state to localStorage."""
    events_data = json.dumps(st.session_state.get("events", []))
    settings_data = json.dumps({
        "timeline_title": st.session_state.get("timeline_title", DEFAULT_TITLE),
        "timeline_view": normalize_view_choice(st.session_state.get("timeline_view", "Timeline")),
        "timeline_bg_color": st.session_state.get("timeline_bg_color", "#fefcea"),
        "event_title_color": st.session_state.get("event_title_color", "#ffffff"),
        "event_title_font": st.session_state.get("event_title_font", "Arial"),
        "event_title_size": st.session_state.get("event_title_size", 12),
        "event_date_color": st.session_state.get("event_date_color", "#ffffff"),
        "event_date_font": st.session_state.get("event_date_font", "Arial"),
        "event_date_size": st.session_state.get("event_date_size", 10),
        "lens_size": st.session_state.get("lens_size", 240),
        "lens_duration": st.session_state.get("lens_duration", 0.5)
    })
    
    js_code = f"""
    <script>
    localStorage.setItem('timeline_events', '{events_data}');
    localStorage.setItem('timeline_settings', '{settings_data}');
    console.log('Data saved to localStorage');
    </script>
    """
    components.html(js_code, height=0)

def load_from_localstorage():
    """Generate JavaScript to load data from localStorage and trigger Streamlit rerun."""
    js_code = """
    <script>
    function loadData() {
        const events = localStorage.getItem('timeline_events');
        const settings = localStorage.getItem('timeline_settings');
        
        if (events || settings) {
            const data = {
                events: events ? JSON.parse(events) : [],
                settings: settings ? JSON.parse(settings) : {}
            };
            
            // Send data to Streamlit
            window.parent.postMessage({
                type: 'streamlit:setComponentValue',
                key: 'loaded_data',
                value: data
            }, '*');
        }
    }
    
    // Try to load data
    loadData();
    
    // Also try loading after a short delay in case localStorage isn't ready
    setTimeout(loadData, 100);
    </script>
    """
    components.html(js_code, height=0)

def clear_localstorage():
    """Generate JavaScript to clear all localStorage data."""
    js_code = """
    <script>
    localStorage.removeItem('timeline_events');
    localStorage.removeItem('timeline_settings');
    console.log('localStorage cleared');
    </script>
    """
    components.html(js_code, height=0)

def create_save_point(save_name):
    """Create a complete save point including events, settings, and images."""
    # Create save points directory if it doesn't exist
    save_dir = Path("save_points")
    save_dir.mkdir(exist_ok=True)
    
    # Prepare save data
    save_data = {
        "save_info": {
            "name": save_name,
            "created_at": datetime.datetime.now().isoformat(),
            "version": "1.0",
            "event_count": len(st.session_state.get("events", []))
        },
        "events": st.session_state.get("events", []),
        "settings": {
            "timeline_title": st.session_state.get("timeline_title", DEFAULT_TITLE),
            "timeline_view": normalize_view_choice(st.session_state.get("timeline_view", "Timeline")),
            "timeline_bg_color": st.session_state.get("timeline_bg_color", "#fefcea"),
            "event_title_color": st.session_state.get("event_title_color", "#ffffff"),
            "event_title_font": st.session_state.get("event_title_font", "Arial"),
            "event_title_size": st.session_state.get("event_title_size", 12),
            "event_date_color": st.session_state.get("event_date_color", "#ffffff"),
            "event_date_font": st.session_state.get("event_date_font", "Arial"),
            "event_date_size": st.session_state.get("event_date_size", 10),
            "lens_size": st.session_state.get("lens_size", 240),
            "lens_duration": st.session_state.get("lens_duration", 0.5)
        }
    }
    
    # Create save point file
    save_file = save_dir / f"{save_name}.json"
    with open(save_file, 'w', encoding='utf-8') as f:
        json.dump(save_data, f, indent=2, ensure_ascii=False)
    
    # Extract and save images
    images_dir = save_dir / f"{save_name}_images"
    images_dir.mkdir(exist_ok=True)
    
    image_files = {}
    for i, event in enumerate(save_data["events"]):
        if event.get("image"):
            # Decode base64 image and save
            try:
                image_data = base64.b64decode(event["image"])
                image_filename = f"event_{i}_{event['id']}.png"
                image_path = images_dir / image_filename
                
                with open(image_path, 'wb') as img_file:
                    img_file.write(image_data)
                
                # Store reference to image file
                image_files[event['id']] = image_filename
                # Remove base64 data from JSON (we'll load from file)
                event["image_file"] = image_filename
                del event["image"]
            except Exception as e:
                st.error(f"Error saving image for event {event['title']}: {e}")
    
    # Update save data with image file references
    save_data["image_files"] = image_files
    
    # Save updated JSON
    with open(save_file, 'w', encoding='utf-8') as f:
        json.dump(save_data, f, indent=2, ensure_ascii=False)
    
    return save_file

def load_save_point(save_name):
    """Load a complete save point including events, settings, and images."""
    save_dir = Path("save_points")
    save_file = save_dir / f"{save_name}.json"
    images_dir = save_dir / f"{save_name}_images"
    
    if not save_file.exists():
        return False, "Save point not found"
    
    try:
        # Load save data
        with open(save_file, 'r', encoding='utf-8') as f:
            save_data = json.load(f)
        
        # Restore settings
        if "settings" in save_data:
            settings = save_data["settings"]
            st.session_state["timeline_title"] = settings.get("timeline_title", DEFAULT_TITLE)
            st.session_state["timeline_view"] = normalize_view_choice(settings.get("timeline_view", "Timeline"))
            st.session_state["timeline_bg_color"] = settings.get("timeline_bg_color", "#fefcea")
            st.session_state["event_title_color"] = settings.get("event_title_color", "#ffffff")
            st.session_state["event_title_font"] = settings.get("event_title_font", "Arial")
            st.session_state["event_title_size"] = settings.get("event_title_size", 12)
            st.session_state["event_date_color"] = settings.get("event_date_color", "#ffffff")
            st.session_state["event_date_font"] = settings.get("event_date_font", "Arial")
            st.session_state["event_date_size"] = settings.get("event_date_size", 10)
            st.session_state["lens_size"] = settings.get("lens_size", 240)
            st.session_state["lens_duration"] = settings.get("lens_duration", 0.5)
        
        # Restore events with images
        events = []
        for event in save_data.get("events", []):
            event_copy = event.copy()
            
            # Load image if available
            if "image_file" in event and images_dir.exists():
                image_path = images_dir / event["image_file"]
                if image_path.exists():
                    try:
                        with open(image_path, 'rb') as img_file:
                            image_data = img_file.read()
                        event_copy["image"] = base64.b64encode(image_data).decode()
                    except Exception as e:
                        st.error(f"Error loading image for event {event['title']}: {e}")
                        event_copy["image"] = None
                
                # Remove temporary image_file reference
                del event_copy["image_file"]
            else:
                event_copy["image"] = None
            
            events.append(event_copy)
        
        st.session_state["events"] = events
        
        return True, f"Successfully loaded save point: {save_name}"
        
    except Exception as e:
        return False, f"Error loading save point: {str(e)}"

def get_save_points():
    """Get list of all available save points."""
    save_dir = Path("save_points")
    if not save_dir.exists():
        return []
    
    save_points = []
    for file_path in save_dir.glob("*.json"):
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                save_data = json.load(f)
            
            save_info = save_data.get("save_info", {})
            save_points.append({
                "name": file_path.stem,
                "created_at": save_info.get("created_at", "Unknown"),
                "event_count": save_info.get("event_count", 0),
                "title": save_info.get("name", file_path.stem)
            })
        except Exception as e:
            st.error(f"Error reading save point {file_path.name}: {e}")
    
    return sorted(save_points, key=lambda x: x["created_at"], reverse=True)

def delete_save_point(save_name):
    """Delete a save point and its associated images."""
    save_dir = Path("save_points")
    save_file = save_dir / f"{save_name}.json"
    images_dir = save_dir / f"{save_name}_images"
    
    deleted_files = []
    
    # Delete save file
    if save_file.exists():
        save_file.unlink()
        deleted_files.append(str(save_file))
    
    # Delete images directory
    if images_dir.exists():
        for img_file in images_dir.glob("*"):
            img_file.unlink()
            deleted_files.append(str(img_file))
        images_dir.rmdir()
        deleted_files.append(str(images_dir))
    
    return deleted_files


@st.cache_resource
def ensure_excel_engine():
    """Ensure an Excel writer engine is available, attempting on-the-fly install if needed."""
    preferred_engines = ("openpyxl", "xlsxwriter")

    for engine in preferred_engines:
        if importlib.util.find_spec(engine) is not None:
            return engine

    missing = [engine for engine in preferred_engines if importlib.util.find_spec(engine) is None]
    if missing:
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", *missing])
        except Exception as exc:
            st.warning(f"Automatic install of Excel export dependencies failed: {exc}")
            return None

    for engine in preferred_engines:
        if importlib.util.find_spec(engine) is not None:
            return engine

    return None

st.set_page_config(
    page_title="Event Timeline Lens",
    page_icon="⏱️",
    layout="wide",
    initial_sidebar_state="expanded",
)

# Custom CSS for modern UI
st.markdown("""
<style>
    /* Main theme styles */
    .stApp {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        background-attachment: fixed;
    }
    
    /* Glassmorphism effect for main content */
    .main .block-container {
        background: rgba(255, 255, 255, 0.1);
        backdrop-filter: blur(10px);
        border-radius: 20px;
        border: 1px solid rgba(255, 255, 255, 0.2);
        padding: 2rem;
        margin-top: 1rem;
    }
    
    /* Sidebar styling */
    .css-1d391kg {
        background: rgba(255, 255, 255, 0.95);
        backdrop-filter: blur(20px);
        border-right: 1px solid rgba(255, 255, 255, 0.3);
    }
    
    /* Title styling */
    .stTitle {
        color: white;
        text-shadow: 2px 2px 4px rgba(0,0,0,0.3);
        font-weight: 700;
        font-size: 2.5rem;
        margin-bottom: 1rem;
    }
    
    /* Button styling */
    .stButton > button {
        background: linear-gradient(45deg, #667eea, #764ba2);
        color: white;
        border: none;
        border-radius: 12px;
        padding: 0.75rem 1.5rem;
        font-weight: 600;
        font-size: 0.95rem;
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
        box-shadow: 0 4px 15px rgba(102, 126, 234, 0.3);
        position: relative;
        overflow: hidden;
        text-transform: none;
        letter-spacing: 0.5px;
    }
    
    .stButton > button::before {
        content: "";
        position: absolute;
        top: 0;
        left: -100%;
        width: 100%;
        height: 100%;
        background: linear-gradient(90deg, transparent, rgba(255, 255, 255, 0.2), transparent);
        transition: left 0.5s;
    }
    
    .stButton > button:hover::before {
        left: 100%;
    }
    
    .stButton > button:hover {
        transform: translateY(-3px) scale(1.02);
        box-shadow: 0 8px 25px rgba(102, 126, 234, 0.4);
        background: linear-gradient(45deg, #764ba2, #667eea);
    }
    
    .stButton > button:active {
        transform: translateY(-1px) scale(0.98);
        box-shadow: 0 4px 15px rgba(102, 126, 234, 0.3);
    }
    
    /* Special styling for primary action buttons */
    .stButton > button[kind="primary"] {
        background: linear-gradient(45deg, #48bb78, #38a169);
        box-shadow: 0 4px 15px rgba(72, 187, 120, 0.3);
    }
    
    .stButton > button[kind="primary"]:hover {
        background: linear-gradient(45deg, #38a169, #2f855a);
        box-shadow: 0 8px 25px rgba(72, 187, 120, 0.4);
    }
    
    /* Expander styling */
    .streamlit-expanderHeader {
        background: linear-gradient(45deg, rgba(102, 126, 234, 0.1), rgba(118, 75, 162, 0.1));
        border-radius: 10px;
        border: 1px solid rgba(255, 255, 255, 0.2);
        font-weight: 600;
        color: #333;
    }
    
    /* Input field styling */
    .stTextInput > div > div > input,
    .stDateInput > div > div > input,
    .stSelectbox > div > div > select {
        background: rgba(255, 255, 255, 0.9);
        border: 1px solid rgba(102, 126, 234, 0.3);
        border-radius: 8px;
        color: #000000;
        transition: all 0.3s ease;
    }
    
    .stTextInput > div > div > input:focus,
    .stDateInput > div > div > input:focus,
    .stSelectbox > div > div > select:focus {
        border-color: #667eea;
        box-shadow: 0 0 0 3px rgba(102, 126, 234, 0.1);
    }
    
    /* Success/info/error message styling */
    .stSuccess, .stInfo, .stWarning, .stError {
        border-radius: 10px;
        border: 1px solid rgba(255, 255, 255, 0.2);
        backdrop-filter: blur(10px);
    }
    
    /* File uploader styling */
    .stFileUploader > div {
        background: rgba(255, 255, 255, 0.1);
        border: 2px dashed rgba(255, 255, 255, 0.3);
        border-radius: 10px;
        transition: all 0.3s ease;
    }
    
    .stFileUploader > div:hover {
        border-color: #667eea;
        background: rgba(102, 126, 234, 0.1);
    }
    
    /* Slider styling */
    .stSlider > div > div > div {
        background: linear-gradient(90deg, #667eea, #764ba2);
    }
    
    /* Color picker styling */
    .stColorPicker > div {
        border-radius: 8px;
        overflow: hidden;
    }
    
    /* Typography improvements */
    h1, h2, h3, h4, h5, h6 {
      font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
      font-weight: 600;
      line-height: 1.3;
    }
    
    p, span, div {
      font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    }
    
    /* Better spacing for sidebar elements */
    .css-1d391kg .stSelectbox, 
    .css-1d391kg .stTextInput, 
    .css-1d391kg .stDateInput,
    .css-1d391kg .stSlider,
    .css-1d391kg .stColorPicker {
      margin-bottom: 1rem;
    }
    
    /* Improved expander spacing */
    .streamlit-expander {
      margin-bottom: 1rem;
    }
    
    .streamlit-expanderContent {
      padding: 1rem;
    }
    
    /* Better column spacing */
    .stColumns > div {
      padding: 0.5rem;
    }
    
    /* Enhanced file upload area */
    .stFileUploader {
      margin: 1rem 0;
    }
    
    /* Better radio button styling */
    .stRadio > div {
      background: rgba(255, 255, 255, 0.05);
      padding: 0.5rem;
      border-radius: 8px;
      border: 1px solid rgba(255, 255, 255, 0.1);
    }
    
    .stRadio option {
      padding: 0.25rem 0;
    }
    
    /* Info message improvements */
    .stInfo {
      background: linear-gradient(45deg, rgba(102, 126, 234, 0.1), rgba(118, 75, 162, 0.1));
      border-left: 4px solid #667eea;
    }
    
    .stSuccess {
      background: linear-gradient(45deg, rgba(72, 187, 120, 0.1), rgba(56, 161, 105, 0.1));
      border-left: 4px solid #48bb78;
    }
    
    .stWarning {
      background: linear-gradient(45deg, rgba(237, 137, 54, 0.1), rgba(221, 107, 32, 0.1));
      border-left: 4px solid #ed8936;
    }
    
    .stError {
      background: linear-gradient(45deg, rgba(245, 101, 101, 0.1), rgba(229, 62, 62, 0.1));
      border-left: 4px solid #f56565;
    }
    
    /* Download button styling */
    .stDownloadButton > button {
      background: linear-gradient(45deg, #48bb78, #38a169);
      color: white;
      border: none;
      border-radius: 12px;
      padding: 0.5rem 1rem;
      font-weight: 600;
      transition: all 0.3s ease;
      box-shadow: 0 4px 15px rgba(0,0,0,0.2);
    }
    
    .stDownloadButton > button:hover {
      transform: translateY(-2px);
      box-shadow: 0 6px 20px rgba(0,0,0,0.3);
      background: linear-gradient(45deg, #38a169, #2f855a);
    }
</style>
""", unsafe_allow_html=True)

if "timeline_title" not in st.session_state:
    st.session_state["timeline_title"] = DEFAULT_TITLE

# Visualization layout preference
if "timeline_view" not in st.session_state:
    st.session_state["timeline_view"] = "Domino"
else:
    st.session_state["timeline_view"] = normalize_view_choice(st.session_state["timeline_view"])

# Initialize session state for events
if "events" not in st.session_state:
    st.session_state["events"] = []

# Try to restore data from localStorage on page load
if "data_restored" not in st.session_state:
    load_from_localstorage()
    st.session_state["data_restored"] = True

# Check for loaded data from localStorage
if "loaded_data" in st.session_state:
    loaded_data = st.session_state["loaded_data"]
    if loaded_data and isinstance(loaded_data, dict):
        # Restore events
        if "events" in loaded_data and loaded_data["events"]:
            st.session_state["events"] = loaded_data["events"]
        
        # Restore settings
        if "settings" in loaded_data and loaded_data["settings"]:
            settings = loaded_data["settings"]
            st.session_state["timeline_title"] = settings.get("timeline_title", DEFAULT_TITLE)
            st.session_state["timeline_view"] = normalize_view_choice(settings.get("timeline_view", "Timeline"))
            st.session_state["timeline_bg_color"] = settings.get("timeline_bg_color", "#fefcea")
            st.session_state["event_title_color"] = settings.get("event_title_color", "#ffffff")
            st.session_state["event_title_font"] = settings.get("event_title_font", "Arial")
            st.session_state["event_title_size"] = settings.get("event_title_size", 12)
            st.session_state["event_date_color"] = settings.get("event_date_color", "#ffffff")
            st.session_state["event_date_font"] = settings.get("event_date_font", "Arial")
            st.session_state["event_date_size"] = settings.get("event_date_size", 10)
            st.session_state["lens_size"] = settings.get("lens_size", 240)
            st.session_state["lens_duration"] = settings.get("lens_duration", 0.5)
    
    # Clear the loaded data to prevent reprocessing
    del st.session_state["loaded_data"]


st.markdown(f"<h1 style='text-align: center; color: white; text-shadow: 2px 2px 4px rgba(0,0,0,0.3); font-weight: 700; font-size: 2.5rem; margin-bottom: 2rem;'>{st.session_state['timeline_title']}</h1>", unsafe_allow_html=True)


# Sidebar for data entry
with st.sidebar:
    st.markdown("""
    <div style="text-align: center; padding: 1rem; margin-bottom: 1rem; background: linear-gradient(45deg, #667eea, #764ba2); border-radius: 15px; color: white;">
        <h2 style="margin: 0; font-size: 1.5rem; font-weight: 700;">⚙️ Control Panel</h2>
        <p style="margin: 0.5rem 0 0; opacity: 0.9; font-size: 0.9rem;">Customize your timeline</p>
    </div>
    """, unsafe_allow_html=True)
    
    # 🎨 Timeline Settings Section
    with st.expander("🎨 Timeline Settings", expanded=True):
        view_options = ("Timeline", "Domino")
        current_view = normalize_view_choice(st.session_state.get("timeline_view", "Timeline"))
        view_choice = st.radio(
            "Visualization Style",
            view_options,
            index=view_options.index(current_view),
            help="Choose between a horizontal timeline or a cascading domino layout."
        )
        st.session_state["timeline_view"] = view_choice

        st.session_state["timeline_title"] = st.text_input(
            "Timeline Title",
            value=st.session_state.get("timeline_title", DEFAULT_TITLE),
            help="Displayed at the top of the timeline."
        ) or DEFAULT_TITLE
        
        # Add spell checking for timeline title
        current_title = st.session_state.get("timeline_title", DEFAULT_TITLE)
        if current_title and current_title != DEFAULT_TITLE and SPELLCHECKER_AVAILABLE:
            spell_checker = get_spell_checker()
            misspelled, suggestions = check_spelling_and_suggest(current_title, spell_checker)
            
            if misspelled:
                st.markdown("🔍 **Timeline Title Spelling:**")
                display_spell_suggestions(current_title, misspelled, suggestions)

        # Background color picker for horizontal timeline
        if "timeline_bg_color" not in st.session_state:
            st.session_state["timeline_bg_color"] = "#fefcea"
        
        timeline_bg_color = st.color_picker(
            "Timeline Background Color",
            value=st.session_state["timeline_bg_color"],
            help="Set the background color for the horizontal timeline view."
        )
        st.session_state["timeline_bg_color"] = timeline_bg_color

    # 🎯 Event Styling Section
    with st.expander("🎯 Event Styling", expanded=False):
        st.markdown("**✨ Event Title Styling**")
        
        # Event title styling options
        if "event_title_color" not in st.session_state:
            st.session_state["event_title_color"] = "#ffffff"
        
        if "event_title_font" not in st.session_state:
            st.session_state["event_title_font"] = "Arial"
        
        if "event_title_size" not in st.session_state:
            st.session_state["event_title_size"] = 12
        
        event_title_color = st.color_picker(
            "🎨 Title Color",
            value=st.session_state["event_title_color"],
            help="Set the color of event titles."
        )
        st.session_state["event_title_color"] = event_title_color
        
        event_title_font = st.selectbox(
            "📝 Title Font",
            options=["Arial", "Helvetica", "Times New Roman", "Georgia", "Verdana", "Courier New", "Impact", "Comic Sans MS"],
            index=["Arial", "Helvetica", "Times New Roman", "Georgia", "Verdana", "Courier New", "Impact", "Comic Sans MS"].index(st.session_state["event_title_font"]),
            help="Choose the font for event titles."
        )
        st.session_state["event_title_font"] = event_title_font
        
        event_title_size = st.slider(
            "📏 Title Size (px)",
            min_value=8,
            max_value=24,
            value=st.session_state["event_title_size"],
            help="Set the font size for event titles in pixels."
        )
        st.session_state["event_title_size"] = event_title_size

        st.markdown("---")
        st.markdown("📅 **Event Date Styling**")
        
        # Event date styling options
        if "event_date_color" not in st.session_state:
            st.session_state["event_date_color"] = "#ffffff"
        
        if "event_date_font" not in st.session_state:
            st.session_state["event_date_font"] = "Arial"
        
        if "event_date_size" not in st.session_state:
            st.session_state["event_date_size"] = 10
        
        event_date_color = st.color_picker(
            "🎨 Date Color",
            value=st.session_state["event_date_color"],
            help="Set the color of event dates."
        )
        st.session_state["event_date_color"] = event_date_color
        
        event_date_font = st.selectbox(
            "📝 Date Font",
            options=["Arial", "Helvetica", "Times New Roman", "Georgia", "Verdana", "Courier New", "Impact", "Comic Sans MS"],
            index=["Arial", "Helvetica", "Times New Roman", "Georgia", "Verdana", "Courier New", "Impact", "Comic Sans MS"].index(st.session_state["event_date_font"]),
            help="Choose the font for event dates."
        )
        st.session_state["event_date_font"] = event_date_font
        
        event_date_size = st.slider(
            "📏 Date Size (px)",
            min_value=6,
            max_value=18,
            value=st.session_state["event_date_size"],
            help="Set the font size for event dates in pixels."
        )
        st.session_state["event_date_size"] = event_date_size

    # 🔍 Lens Settings Section
    with st.expander("🔍 Lens Settings", expanded=False):
        st.markdown("**🔍 Magnification Settings**")
        
        # Lens styling options
        if "lens_size" not in st.session_state:
            st.session_state["lens_size"] = 240
        
        lens_size = st.slider(
            "🔍 Lens Size (px)",
            min_value=120,
            max_value=400,
            value=st.session_state["lens_size"],
            help="Set the size of the lens in pixels."
        )
        st.session_state["lens_size"] = lens_size

        # Lens magnification duration
        if "lens_duration" not in st.session_state:
            st.session_state["lens_duration"] = 0.5
        
        lens_duration = st.slider(
            "⏱️ Magnification Delay (seconds)",
            min_value=0.0,
            max_value=3.0,
            value=st.session_state["lens_duration"],
            step=0.1,
            help="Delay before magnification effect when event enters the lens."
        )
        st.session_state["lens_duration"] = lens_duration

    # ➕ Add Event Section
    with st.expander("➕ Add New Event", expanded=True):
        st.markdown("**📝 Create New Event**")
        
        title_input = st.text_input(
            "🎯 Event Title",
            placeholder="Enter event title...",
            help="Give your event a descriptive title"
        )
        
        # Add spell checking for event title
        if title_input and SPELLCHECKER_AVAILABLE:
            spell_checker = get_spell_checker()
            misspelled, suggestions = check_spelling_and_suggest(title_input, spell_checker)
            
            with st.expander("🔍 Check Spelling", expanded=False):
                display_spell_suggestions(title_input, misspelled, suggestions)
        
        date_input = st.date_input(
            "📅 Event Date", 
            datetime.date(2002, 1, 1), 
            min_value=datetime.date(2002, 1, 1), 
            max_value=datetime.date.today(),
            help="Select the date for your event"
        )
        
        image_file = st.file_uploader(
            "🖼️ Upload Image (optional)", 
            type=["png", "jpg", "jpeg", "gif"],
            help="Add an image to make your event more visual"
        )
        image_data = None
        if image_file:
            image_data = base64.b64encode(image_file.read()).decode()

        if st.button("✨ Add Event", help="Click to add this event to your timeline"):
            if title_input:
                st.session_state["events"].append(
                    {
                        "id": str(uuid4()),
                        "title": title_input,
                        "date": date_input.isoformat(),
                        "image": image_data,
                    }
                )
                # Auto-save to localStorage
                save_to_localstorage()
                # Force a rerun so the timeline updates immediately
                rerun_app()
            else:
                st.warning("Please enter an event title before adding.")

    # 📊 Import/Export Section
    with st.expander("📊 Import & Export", expanded=False):
        st.markdown("**Import Events from Excel**")
        st.info("Excel file must contain columns named 'eventname' and 'eventdate' (case-insensitive)")
        excel_file = st.file_uploader("Upload Excel file", type=["xlsx", "xls", "xlsb"], key="excel_uploader")
        if excel_file:
            try:
                ext = excel_file.name.split(".")[-1].lower()
                read_kwargs = {}
                if ext == "xlsb":
                    if importlib.util.find_spec("pyxlsb") is None:
                        st.error("pyxlsb package is required for .xlsb files. Please install it with: pip install pyxlsb")
                        st.stop()
                    read_kwargs["engine"] = "pyxlsb"
                
                st.info(f"Reading Excel file: {excel_file.name}")
                df = pd.read_excel(excel_file, **read_kwargs)
                
                # Show file info
                st.write(f"Found {len(df)} rows in Excel file")
                st.write("Columns found:", df.columns.tolist())
                
                missing_cols = {"eventname", "eventdate"} - set(col.lower() for col in df.columns)
                if missing_cols:
                    st.error(f"Excel file must contain columns: eventname, eventdate. Missing: {missing_cols}")
                    st.write("Available columns:", [col for col in df.columns])
                else:
                    normalized = {col.lower(): col for col in df.columns}
                    new_events = []
                    skipped_count = 0
                    duplicate_count = 0
                    
                    # Get existing events for duplicate checking
                    existing_events = set()
                    for ev in st.session_state["events"]:
                        key = (ev["title"].lower().strip(), ev["date"])
                        existing_events.add(key)
                    
                    for idx, row in df.iterrows():
                        name = row[normalized["eventname"]]
                        date_val = row[normalized["eventdate"]]
                        
                        if pd.isna(name) or pd.isna(date_val):
                            skipped_count += 1
                            continue
                            
                        try:
                            # Handle pandas Timestamp objects (most common from Excel)
                            if hasattr(date_val, 'date'):
                                parsed_date = date_val.date()
                            # Handle datetime objects
                            elif isinstance(date_val, datetime.datetime):
                                parsed_date = date_val.date()
                            # Handle date objects
                            elif isinstance(date_val, datetime.date):
                                parsed_date = date_val
                            # Handle string dates
                            else:
                                parsed_date = pd.to_datetime(date_val).date()
                            
                            # Validate minimum date constraint
                            if parsed_date < datetime.date(2002, 1, 1):
                                st.warning(f"Row {idx+1}: Skipping event '{name}' - date {parsed_date} is before minimum allowed date (2002-01-01)")
                                skipped_count += 1
                                continue
                                
                            date_iso = parsed_date.isoformat()
                        except Exception as e:
                            st.warning(f"Row {idx+1}: Skipping event '{name}' - invalid date format: {date_val}")
                            skipped_count += 1
                            continue
                        
                        # Check for duplicates
                        event_key = (str(name).lower().strip(), date_iso)
                        if event_key in existing_events:
                            duplicate_count += 1
                            continue
                            
                        existing_events.add(event_key)
                            
                        new_events.append(
                            {
                                "id": str(uuid4()),
                                "title": str(name).strip(),
                                "date": date_iso,
                                "image": None,
                            }
                        )
                    
                    if new_events:
                        st.session_state["events"].extend(new_events)
                        st.success(f"Successfully imported {len(new_events)} events from Excel.")
                        if skipped_count > 0:
                            st.info(f"Skipped {skipped_count} rows due to missing data or invalid dates")
                        if duplicate_count > 0:
                            st.info(f"Skipped {duplicate_count} duplicate events that already exist")
                        # Clear the file uploader to prevent re-importing on rerun
                        st.session_state["excel_uploader"] = None
                        st.rerun()
                    else:
                        st.warning("No valid rows found to import. Check that your Excel file has proper eventname and eventdate columns.")
            except Exception as exc:
                st.error(f"Failed to read Excel file: {exc}")
                st.write("Please ensure your Excel file is not corrupted and has the correct format.")

        st.markdown("---")
        st.markdown("**Export Events to Excel**")
        if st.session_state["events"]:
            sorted_export = sorted(
                st.session_state["events"],
                key=lambda e: datetime.date.fromisoformat(e["date"]),
            )
            export_rows = [
                {
                    "eventname": ev["title"],
                    "eventdate": ev["date"],
                    "image": ev.get("image"),
                }
                for ev in sorted_export
            ]
            export_df = pd.DataFrame(export_rows)
            excel_buffer = BytesIO()

            engine = ensure_excel_engine()
            if engine is None:
                st.error(
                    "Excel export requires the 'openpyxl' or 'XlsxWriter' package. "
                    "Please install one of them in the deployment environment."
                )
                excel_buffer = None
            else:
                try:
                    export_df.to_excel(excel_buffer, index=False, engine=engine)
                except ModuleNotFoundError:
                    st.error(
                        "Excel export requires the 'openpyxl' or 'XlsxWriter' package. "
                        "Please install one of them in the deployment environment."
                    )
                    excel_buffer = None

            if excel_buffer is not None:
                st.download_button(
                    label="Download events.xlsx",
                    data=excel_buffer.getvalue(),
                    file_name="timeline_events.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
        else:
            st.info("Add events to enable exporting.")

        st.markdown("---")
        st.markdown("**📄 Export as PDF**")
        if st.session_state["events"]:
            st.info("Generate a PDF document with your timeline events and images")
            
            if st.button("📄 Generate PDF"):
                try:
                    # Import required libraries
                    from reportlab.lib.pagesizes import letter, A4
                    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
                    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                    from reportlab.lib.units import inch
                    from reportlab.lib import colors
                    from io import BytesIO
                    import base64
                    
                    # Create PDF buffer
                    pdf_buffer = BytesIO()
                    doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
                    
                    # Get styles
                    styles = getSampleStyleSheet()
                    title_style = ParagraphStyle(
                        'CustomTitle',
                        parent=styles['Heading1'],
                        fontSize=24,
                        spaceAfter=30,
                        textColor=colors.darkblue,
                        alignment=1  # Center
                    )
                    
                    heading_style = ParagraphStyle(
                        'CustomHeading',
                        parent=styles['Heading2'],
                        fontSize=16,
                        spaceAfter=12,
                        textColor=colors.darkblue
                    )
                    
                    # Build PDF content
                    story = []
                    
                    # Title
                    story.append(Paragraph(st.session_state.get("timeline_title", "Event Timeline"), title_style))
                    story.append(Spacer(1, 20))
                    
                    # Events
                    sorted_events = sorted(
                        st.session_state["events"],
                        key=lambda e: datetime.date.fromisoformat(e["date"])
                    )
                    
                    for i, event in enumerate(sorted_events):
                        # Event heading
                        event_title = f"{event['title']} - {event['date']}"
                        story.append(Paragraph(event_title, heading_style))
                        
                        # Event image if available
                        if event.get("image"):
                            try:
                                # Decode base64 image
                                image_data = base64.b64decode(event['image'])
                                image_buffer = BytesIO(image_data)
                                
                                # Add image to PDF (scaled to fit)
                                img = Image(image_buffer, width=4*inch, height=3*inch)
                                story.append(img)
                                story.append(Spacer(1, 12))
                            except Exception as img_err:
                                st.warning(f"Could not include image for event: {event['title']}")
                        
                        story.append(Spacer(1, 20))
                    
                    # Build PDF
                    doc.build(story)
                    
                    # Provide download
                    pdf_buffer.seek(0)
                    st.download_button(
                        label="📄 Download Timeline PDF",
                        data=pdf_buffer.getvalue(),
                        file_name="timeline_events.pdf",
                        mime="application/pdf"
                    )
                    
                except ImportError:
                    st.error("""
                    **PDF export requires additional packages.** 
                    Please install: `pip install reportlab`
                    """)
                except Exception as e:
                    st.error(f"Error generating PDF: {str(e)}")
        else:
            st.info("Add events to enable PDF export.")

        st.markdown("---")
        st.markdown("**📊 Export as PowerPoint**")
        if st.session_state["events"]:
            st.info("Create a PowerPoint presentation with your timeline events")
            
            if st.button("📊 Generate PowerPoint"):
                try:
                    # Import required libraries
                    from pptx import Presentation
                    from pptx.util import Inches, Pt
                    from pptx.enum.text import PP_ALIGN
                    from pptx.dml.color import RGBColor
                    from io import BytesIO
                    import base64
                    
                    # Create presentation
                    prs = Presentation()
                    
                    # Title slide
                    title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
                    title = title_slide.shapes.title
                    subtitle = title_slide.placeholders[1]
                    
                    title.text = st.session_state.get("timeline_title", "Event Timeline")
                    subtitle.text = f"Generated on {datetime.date.today().strftime('%B %d, %Y')} • {len(st.session_state['events'])} events"
                    
                    # Event slides
                    sorted_events = sorted(
                        st.session_state["events"],
                        key=lambda e: datetime.date.fromisoformat(e["date"])
                    )
                    
                    for event in sorted_events:
                        # Use title and content layout
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        
                        # Set title
                        title_shape = slide.shapes.title
                        title_shape.text = f"{event['title']}"
                        
                        # Add date as subtitle
                        content_placeholder = slide.placeholders[1]
                        content_placeholder.text = f"📅 {event['date']}"
                        
                        # Add image if available
                        if event.get("image"):
                            try:
                                # Decode base64 image
                                image_data = base64.b64decode(event['image'])
                                image_buffer = BytesIO(image_data)
                                
                                # Add image to slide (positioned below title)
                                left = Inches(1)
                                top = Inches(2.5)
                                width = Inches(8)
                                height = Inches(4)
                                
                                slide.shapes.add_picture(image_buffer, left, top, width=width, height=height)
                                
                                # Adjust text position if image added
                                content_placeholder.top = Inches(6.5)
                                
                            except Exception as img_err:
                                st.warning(f"Could not include image for event: {event['title']}")
                    
                    # Save presentation to buffer
                    pptx_buffer = BytesIO()
                    prs.save(pptx_buffer)
                    
                    # Provide download
                    pptx_buffer.seek(0)
                    st.download_button(
                        label="📊 Download PowerPoint Presentation",
                        data=pptx_buffer.getvalue(),
                        file_name="timeline_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    
                except ImportError:
                    st.error("""
                    **PowerPoint export requires additional packages.** 
                    Please install: `pip install python-pptx`
                    """)
                except Exception as e:
                    st.error(f"Error generating PowerPoint: {str(e)}")
        else:
            st.info("Add events to enable PowerPoint export.")

# Editing existing events
with st.sidebar:
    # ✏️ Edit Events Section (no outer expander to avoid nesting)
    st.markdown("### ✏️ Edit Previous Events")
    for idx, ev in enumerate(st.session_state["events"]):
        with st.expander(f"{ev['title']} - {ev['date']}", expanded=False):
            new_title = st.text_input(f"Title_{ev['id']}", ev["title"], key=f"title_{ev['id']}_edit")
            
            # Add spell checking for edited event title
            if new_title and new_title != ev["title"] and SPELLCHECKER_AVAILABLE:
                spell_checker = get_spell_checker()
                misspelled, suggestions = check_spelling_and_suggest(new_title, spell_checker)
                
                if misspelled:
                    st.markdown(f"🔍 **Spelling Check for '{ev['title']}':**")
                    display_spell_suggestions(new_title, misspelled, suggestions)
            
            new_date = st.date_input(
                f"Date_{ev['id']}",
                datetime.date.fromisoformat(ev['date']),
                key=f"date_{ev['id']}_edit",
                min_value=datetime.date(2002, 1, 1),
                max_value=datetime.date.today()
            )
            new_image_file = st.file_uploader(
                "Replace image (optional)",
                type=["png", "jpg", "jpeg", "gif"],
                key=f"image_{ev['id']}_edit"
            )
            new_image_data = ev.get('image')
            if new_image_file:
                new_image_data = base64.b64encode(new_image_file.read()).decode()
            col1, col2 = st.columns(2)
            if col1.button("Save", key=f"save_{ev['id']}"):
                st.session_state['events'][idx]['title'] = new_title
                st.session_state['events'][idx]['date'] = new_date.isoformat()
                st.session_state['events'][idx]['image'] = new_image_data
                # Auto-save to localStorage
                save_to_localstorage()
                rerun_app()
            if col2.button("Delete", key=f"delete_{ev['id']}"):
                st.session_state['events'].pop(idx)
                # Auto-save to localStorage
                save_to_localstorage()
                rerun_app()

    # 💾 Data Management Section
    with st.expander("💾 Data Management", expanded=False):
        st.markdown("**Browser Storage (Auto-save)**")
        col1, col2, col3 = st.columns(3)
        
        with col1:
            if st.button("💾 Save Now", help="Manually save all events and settings to browser storage"):
                save_to_localstorage()
                st.success("Data saved successfully!")
                st.rerun()
        
        with col2:
            if st.button("🔄 Restore", help="Restore data from browser storage"):
                load_from_localstorage()
                st.info("Attempting to restore data...")
                st.rerun()
        
        with col3:
            if st.button("🗑️ Clear All", help="Clear all saved data from browser storage"):
                clear_localstorage()
                st.session_state["events"] = []
                st.session_state["timeline_title"] = DEFAULT_TITLE
                st.session_state["timeline_view"] = "Timeline"
                st.warning("All data cleared!")
                st.rerun()
        
        st.info("💡 **Auto-save enabled**: Your events and settings are automatically saved when you add, edit, or delete items. Data persists across page refreshes!")

    # 📁 Save Points Section (no outer expander to avoid nesting)
    st.markdown("### 📁 Save Points (Permanent Backup)")
    st.info("📁 **Save Points** are permanent backups stored as files on your computer. They include all events, settings, and uploaded images. You can create multiple save points and restore them anytime, even after closing the app.")
    
    # Create new save point
    with st.expander("🆕 Create New Save Point", expanded=False):
        save_name = st.text_input("Save Point Name", placeholder="e.g., My Timeline v1, Project Timeline, Backup 2024-01-15")
        
        if st.button("💾 Create Save Point", disabled=not save_name.strip()):
            save_name = save_name.strip()
            try:
                save_file = create_save_point(save_name)
                st.success(f"✅ Save point '{save_name}' created successfully!")
                st.info(f"📁 Saved to: {save_file}")
                st.rerun()
            except Exception as e:
                st.error(f"❌ Error creating save point: {str(e)}")
    
    # List and manage existing save points
    save_points = get_save_points()
    
    if save_points:
        st.write(f"📋 **{len(save_points)} Save Points Available:**")
        
        for save_point in save_points:
            with st.expander(f"📁 {save_point['title']} ({save_point['event_count']} events)", expanded=False):
                col1, col2, col3 = st.columns([2, 1, 1])
                
                with col1:
                    st.write(f"**Created:** {save_point['created_at']}")
                    st.write(f"**Events:** {save_point['event_count']}")
                    st.write(f"**Files:** `save_points/{save_point['name']}.json` + images")
                
                with col2:
                    if st.button("🔄 Load", key=f"load_{save_point['name']}", help="Load this save point"):
                        success, message = load_save_point(save_point['name'])
                        if success:
                            st.success(message)
                            # Also save to localStorage for persistence
                            save_to_localstorage()
                            st.rerun()
                        else:
                            st.error(message)
                
                with col3:
                    if st.button("🗑️ Delete", key=f"delete_{save_point['name']}", help="Delete this save point permanently"):
                        if st.session_state.get(f"confirm_delete_{save_point['name']}", False):
                            deleted_files = delete_save_point(save_point['name'])
                            st.success(f"✅ Save point '{save_point['name']}' deleted!")
                            st.info(f"🗑️ Deleted files: {len(deleted_files)}")
                            st.session_state[f"confirm_delete_{save_point['name']}"] = False
                            st.rerun()
                        else:
                            st.session_state[f"confirm_delete_{save_point['name']}"] = True
                            st.warning("⚠️ Click again to confirm deletion")
    else:
        st.info("📝 No save points found. Create your first save point above to permanently backup your timeline!")

# If no events, prompt the user and stop
if not st.session_state["events"]:
    st.info("Use the sidebar to add events. They will appear in a colorful, animated timeline!")
    st.stop()

# Sort events chronologically (oldest first)
sorted_events = sorted(
    st.session_state["events"],
    key=lambda e: datetime.date.fromisoformat(e["date"]),
)

# Serialize to JSON for JavaScript
events_json = json.dumps(sorted_events)
view_mode = st.session_state["timeline_view"]

# Show warning for too many events in domino mode
if view_mode == "Domino" and len(sorted_events) > 40:
    st.warning(f"⚠️ You have {len(sorted_events)} events in domino mode. For better visibility, consider switching to Timeline view or reducing the number of events to under 40.")

# HTML + CSS + JavaScript for the timeline
domino_html = f"""
<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <style>
    html, body {{
      height: 100%;
    }}
    :root {{
      --title-color: {st.session_state["event_title_color"]};
      --date-color: {st.session_state["event_date_color"]};
      --title-font: '{st.session_state["event_title_font"]}';
      --date-font: '{st.session_state["event_date_font"]}';
    }}
    body {{
      margin: 0;
      background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
      font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
      color: #ffffff;
      height: 100%;
      overflow: hidden;
    }}
    .domino-stage {{
      padding: 24px clamp(24px, 6vw, 120px) 32px;
      position: relative;
      height: 100%;
      box-sizing: border-box;
      display: flex;
      flex-direction: column;
      gap: 18px;
      background: radial-gradient(circle at 30% 20%, rgba(255, 255, 255, 0.1) 0%, transparent 50%),
                  radial-gradient(circle at 70% 80%, rgba(255, 255, 255, 0.08) 0%, transparent 50%);
    }}
    .domino-headline {{
      display: flex;
      flex-wrap: wrap;
      align-items: center;
      justify-content: space-between;
      gap: 16px;
      margin-bottom: 24px;
    }}
    .domino-headline h2 {{
      margin: 0;
      font-size: clamp(1.4rem, 2vw, 2rem);
      letter-spacing: 0.04em;
      background: linear-gradient(45deg, #ffffff, #e0e0ff);
      -webkit-background-clip: text;
      -webkit-text-fill-color: transparent;
      background-clip: text;
      text-shadow: 0 0 30px rgba(255, 255, 255, 0.3);
    }}
    .domino-headline p {{
      margin: 4px 0 0;
      opacity: 0.8;
    }}
    #replay-domino {{
      background: rgba(255, 255, 255, 0.15);
      color: #ffffff;
      border: 1px solid rgba(255, 255, 255, 0.35);
      border-radius: 999px;
      padding: 10px 18px;
      cursor: pointer;
      font-weight: 600;
      transition: background 0.3s ease, transform 0.2s ease, border-color 0.3s ease;
    }}
    #replay-domino:disabled {{
      opacity: 0.4;
      cursor: not-allowed;
    }}
    #replay-domino:not(:disabled):hover {{
      background: rgba(255, 255, 255, 0.3);
      transform: translateY(-2px);
      border-color: rgba(255, 255, 255, 0.6);
    }}
    .domino-track {{
      display: flex;
      gap: clamp(12px, 2vw, 28px);
      overflow-x: auto;
      padding: 20px 0 24px 0;
      scroll-snap-type: x proximity;
      flex: 1;
      align-items: stretch;
      min-height: 60px;
    }}
    .domino-track::-webkit-scrollbar {{
      height: 6px;
    }}
    .domino-track::-webkit-scrollbar-thumb {{
      background: rgba(255, 255, 255, 0.25);
      border-radius: 999px;
    }}
    .domino {{
      width: clamp(120px, 12vw, 180px);
      min-height: 40px;
      border-radius: 18px;
      padding: 8px;
      background: linear-gradient(155deg, rgba(255, 255, 255, 0.25), rgba(255, 255, 255, 0.05));
      border: 2px solid rgba(255, 255, 255, 0.35);
      box-shadow: 0 20px 35px rgba(0, 0, 0, 0.4), 0 0 15px rgba(102, 126, 234, 0.3);
      transform-origin: bottom center;
      transform: perspective(900px) rotateX(10deg) rotateY(-4deg) rotateZ(0deg);
      transition: transform 0.45s cubic-bezier(0.4, 0, 0.2, 1), box-shadow 0.45s ease, border-color 0.45s ease, background 0.45s ease;
      cursor: pointer;
      position: relative;
      opacity: 0;
      animation-name: dominoEnter;
      animation-duration: 0.85s;
      animation-timing-function: cubic-bezier(0.4, 0, 0.2, 1);
      animation-fill-mode: forwards;
      animation-delay: var(--delay, 0ms);
      scroll-snap-align: start;
    }}
    .domino::before {{
      content: "";
      position: absolute;
      inset: -4px;
      border-radius: 24px;
      background: radial-gradient(circle at 30% 30%, rgba(255, 255, 255, 0.35), rgba(255, 255, 255, 0));
      opacity: 0;
      transition: opacity 0.45s ease;
      pointer-events: none;
      z-index: -1;
    }}
    .domino::after {{
      content: "";
      position: absolute;
      inset: 12px;
      border-radius: 18px;
      border: 1px dashed rgba(255, 255, 255, 0.18);
      pointer-events: none;
      opacity: 0.4;
    }}
    .domino-image {{
      width: 100%;
      height: clamp(20px, 4vw, 30px);
      border-radius: 8px;
      overflow: hidden;
      margin-bottom: 4px;
      background: linear-gradient(45deg, rgba(0, 0, 0, 0.3), rgba(0, 0, 0, 0.1));
      display: flex;
      align-items: center;
      justify-content: center;
      position: relative;
      border: 1px solid rgba(255, 255, 255, 0.1);
    }}
    .domino-image img {{
      width: 100%;
      height: 100%;
      object-fit: cover;
      filter: saturate(0.8);
      transition: filter 0.35s ease;
    }}
    .domino-placeholder {{
      font-size: clamp(0.5rem, 1.2vw, 0.8rem);
      font-weight: 600;
      text-transform: none;
      opacity: 0.9;
      letter-spacing: 0.02em;
      text-align: center;
      color: #ffffff;
      padding: 2px;
      background: linear-gradient(45deg, #667eea, #764ba2);
      -webkit-background-clip: text;
      -webkit-text-fill-color: transparent;
      background-clip: text;
      overflow: hidden;
      text-overflow: ellipsis;
      white-space: nowrap;
      width: 100%;
    }}
    .domino-date {{
      font-size: clamp({st.session_state["event_date_size"] - 6}px, 1.5vw, {st.session_state["event_date_size"] - 2}px);
      color: var(--date-color);
      font-family: var(--date-font);
      letter-spacing: 0.08em;
      margin: 1px 0;
    }}
    .domino-title {{
      font-size: clamp({st.session_state["event_title_size"] - 8}px, 2vw, {st.session_state["event_title_size"] - 4}px);
      color: var(--title-color);
      font-family: var(--title-font);
      margin: 1px 0 0;
      min-height: 12px;
      line-height: 1.0;
    }}
    .domino:hover,
    .domino.hovering,
    .domino.pinned {{
      transform: perspective(900px) rotateX(4deg) rotateY(0deg) rotateZ(-6deg) translateY(-8px);
      border-color: rgba(255, 255, 255, 0.9);
      box-shadow: 0 20px 40px rgba(0, 0, 0, 0.6), 0 0 20px rgba(102, 126, 234, 0.6), 0 0 30px rgba(118, 75, 162, 0.3);
      background: linear-gradient(150deg, rgba(255, 255, 255, 0.45), rgba(255, 255, 255, 0.15));
    }}
    .domino:hover::before,
    .domino.hovering::before,
    .domino.pinned::before {{
      opacity: 0.8;
      background: radial-gradient(circle at 30% 30%, rgba(255, 255, 255, 0.6), rgba(255, 255, 255, 0));
    }}
    .domino.pinned {{
      border-color: rgba(255, 255, 255, 1);
      box-shadow: 0 25px 50px rgba(0, 0, 0, 0.7), 0 0 30px rgba(102, 126, 234, 0.7), 0 0 40px rgba(118, 75, 162, 0.4);
    }}
    .domino.is-tilting {{
      animation: dominoTip 0.9s cubic-bezier(0.68, -0.2, 0.265, 1.2);
    }}
    @keyframes dominoEnter {{
      0% {{
        opacity: 0;
        transform: perspective(900px) rotateX(60deg) rotateY(-20deg) rotateZ(18deg) translateY(35px);
      }}
      100% {{
        opacity: 1;
        transform: perspective(900px) rotateX(10deg) rotateY(-4deg) rotateZ(0deg);
      }}
    }}
    @keyframes dominoTip {{
      0% {{
        transform: perspective(900px) rotateX(10deg) rotateY(-4deg) rotateZ(0deg);
      }}
      70% {{
        transform: perspective(900px) rotateX(14deg) rotateY(6deg) rotateZ(-16deg) translateY(-8px);
      }}
      100% {{
        transform: perspective(900px) rotateX(6deg) rotateY(0deg) rotateZ(-10deg);
      }}
    }}
    .domino-detail {{
      margin-top: 10px;
      padding: clamp(20px, 3vw, 36px);
      border-radius: 28px;
      border: 1px solid rgba(255, 255, 255, 0.3);
      background: linear-gradient(135deg, rgba(0, 0, 0, 0.5), rgba(0, 0, 0, 0.25)),
                  radial-gradient(circle at 70% 30%, rgba(102, 126, 234, 0.1), transparent 50%);
      backdrop-filter: blur(25px);
      display: flex;
      flex-wrap: wrap;
      gap: 24px;
      align-items: center;
      transition: opacity 0.45s cubic-bezier(0.4, 0, 0.2, 1), transform 0.45s cubic-bezier(0.4, 0, 0.2, 1);
      flex-shrink: 0;
      box-shadow: 0 15px 40px rgba(0, 0, 0, 0.4), 0 0 20px rgba(102, 126, 234, 0.2);
      position: relative;
      overflow: hidden;
    }}
    .domino-detail.hidden {{
      opacity: 0;
      transform: translateY(30px) scale(0.95);
      pointer-events: none;
    }}
    .detail-media {{
      flex: 0 1 clamp(200px, 30vw, 320px);
      height: clamp(160px, 25vw, 220px);
      border-radius: 22px;
      overflow: hidden;
      background: linear-gradient(45deg, rgba(255, 255, 255, 0.1), rgba(255, 255, 255, 0.05));
      display: flex;
      align-items: center;
      justify-content: center;
      border: 1px solid rgba(255, 255, 255, 0.15);
    }}
    .detail-media img {{
      width: 100%;
      height: 100%;
      object-fit: cover;
    }}
    .detail-placeholder {{
      font-size: clamp(1rem, 2.5vw, 1.5rem);
      font-weight: 600;
      opacity: 0.9;
      letter-spacing: 0.02em;
      text-align: center;
      text-transform: none;
      background: linear-gradient(45deg, #667eea, #764ba2);
      -webkit-background-clip: text;
      -webkit-text-fill-color: transparent;
      background-clip: text;
      padding: 10px;
      width: 100%;
      height: 100%;
      display: flex;
      align-items: center;
      justify-content: center;
    }}
    .detail-text {{
      flex: 1 1 260px;
    }}
    .detail-text h3 {{
      margin: 0 0 8px;
      font-size: clamp(1.3rem, 2.2vw, 2rem);
      color: var(--title-color);
      font-family: var(--title-font);
    }}
    .detail-text p {{
      margin: 4px 0;
      color: var(--date-color);
      font-family: var(--date-font);
      letter-spacing: 0.05em;
    }}
    .detail-hint {{
      opacity: 0.65;
      color: #f5f5f5;
    }}
    @media (max-width: 768px) {{
      .domino-stage {{
        padding: 24px 16px 60px;
      }}
      .domino-track {{
        gap: 12px;
      }}
    }}
  </style>
</head>
<body>
  <div class="domino-stage">
    <div class="domino-headline">
      <div>
        <h2>Domino Effect Timeline</h2>
        <p id="domino-status">Preparing your events...</p>
      </div>
      <button id="replay-domino" type="button">Replay Effect</button>
    </div>
    <div id="domino-track" class="domino-track"></div>
    <section id="domino-detail" class="domino-detail hidden">
      <div id="detail-media" class="detail-media"></div>
      <div class="detail-text">
        <h3 id="detail-title"></h3>
        <p id="detail-date"></p>
        <p class="detail-hint">Hover to pin an event, move mouse away to unpin.</p>
      </div>
    </section>
  </div>
  <script>
    const events = {events_json};
    const dataset = Array.isArray(events) ? [...events] : [];
    dataset.sort((a, b) => new Date(a.date) - new Date(b.date));
    const track = document.getElementById('domino-track');
    const detailBlock = document.getElementById('domino-detail');
    const detailTitle = document.getElementById('detail-title');
    const detailDate = document.getElementById('detail-date');
    const detailMedia = document.getElementById('detail-media');
    const statusLabel = document.getElementById('domino-status');
    const replayButton = document.getElementById('replay-domino');
    let dominoElements = [];
    let activeDomino = null;

    function formatDate(value) {{
      if (!value) {{
        return 'No date';
      }}
      const parsed = new Date(value);
      if (Number.isNaN(parsed.valueOf())) {{
        return value;
      }}
      return parsed.toLocaleDateString(undefined, {{ year: 'numeric', month: 'short', day: 'numeric' }});
    }}

    function createPlaceholderChar(text) {{
      if (!text || typeof text !== 'string') {{
        return '#';
      }}
      const trimmed = text.trim();
      return trimmed ? trimmed.charAt(0).toUpperCase() : '#';
    }}

    function buildDominos() {{
      track.innerHTML = '';
      dominoElements = [];
      dataset.forEach((event, index) => {{
        const card = document.createElement('article');
        card.className = 'domino';
        card.style.setProperty('--delay', `${{index * 110}}ms`);
        card.dataset.index = index;

        const mediaWrapper = document.createElement('div');
        mediaWrapper.className = 'domino-image';
        if (event.image) {{
          const img = document.createElement('img');
          img.src = `data:image/png;base64,${{event.image}}`;
          img.alt = event.title || 'Event image';
          mediaWrapper.appendChild(img);
        }} else {{
          const placeholder = document.createElement('span');
          placeholder.className = 'domino-placeholder';
          placeholder.textContent = event.title || 'Event';
          mediaWrapper.appendChild(placeholder);
        }}
        card.appendChild(mediaWrapper);

        const dateEl = document.createElement('p');
        dateEl.className = 'domino-date';
        dateEl.textContent = formatDate(event.date);
        card.appendChild(dateEl);

        const titleEl = document.createElement('h3');
        titleEl.className = 'domino-title';
        titleEl.textContent = event.title || 'Untitled Event';
        card.appendChild(titleEl);

        card.addEventListener('mouseenter', () => {{
          pinDomino(card, event);
        }});
        card.addEventListener('mouseleave', () => {{
          if (card === activeDomino) {{
            card.classList.remove('pinned');
            activeDomino = null;
            detailBlock.classList.add('hidden');
          }}
        }});
        card.addEventListener('focus', () => pinDomino(card, event));
        card.addEventListener('blur', () => {{
          if (card === activeDomino) {{
            card.classList.remove('pinned');
            activeDomino = null;
            detailBlock.classList.add('hidden');
          }}
        }});
        card.addEventListener('click', () => pinDomino(card, event));

        track.appendChild(card);
        dominoElements.push(card);
      }});

      if (!dominoElements.length) {{
        statusLabel.textContent = 'No events yet — add some from the sidebar.';
        replayButton.disabled = true;
        detailBlock.classList.add('hidden');
        return;
      }}

      replayButton.disabled = false;
      statusLabel.textContent = 'Hover to preview, click to pin an event.';
      pinDomino(dominoElements[0], dataset[0]);
      triggerDominoChain();
    }}

    function pinDomino(domino, eventData) {{
      if (!domino || !eventData) {{
        return;
      }}
      dominoElements.forEach(el => el.classList.remove('pinned'));
      domino.classList.add('pinned');
      activeDomino = domino;
      updateDetail(eventData);
      domino.scrollIntoView({{ behavior: 'smooth', inline: 'center', block: 'nearest' }});
    }}

    function updateDetail(eventData) {{
      detailTitle.textContent = eventData.title || 'Untitled Event';
      detailDate.textContent = formatDate(eventData.date);
      detailMedia.innerHTML = '';
      if (eventData.image) {{
        const img = document.createElement('img');
        img.src = `data:image/png;base64,${{eventData.image}}`;
        img.alt = eventData.title || 'Event image';
        detailMedia.appendChild(img);
      }} else {{
        const placeholder = document.createElement('span');
        placeholder.className = 'detail-placeholder';
        placeholder.textContent = eventData.title || 'Event';
        detailMedia.appendChild(placeholder);
      }}
      detailBlock.classList.remove('hidden');
    }}

    function triggerDominoChain() {{
      dominoElements.forEach((domino, index) => {{
        setTimeout(() => {{
          domino.classList.add('is-tilting');
          setTimeout(() => domino.classList.remove('is-tilting'), 950);
        }}, index * 150);
      }});
    }}

    replayButton.addEventListener('click', () => {{
      // Reset all dominoes
      dominoElements.forEach(domino => {{
        domino.classList.remove('pinned', 'hovering');
        domino.style.setProperty('--delay', '0ms');
      }});
      
      // Clear any existing detail
      const existingDetail = document.querySelector('.domino-detail');
      if (existingDetail) {{
        existingDetail.classList.add('hidden');
      }}
      
      // Trigger the chain reaction
      triggerDominoChain();
      statusLabel.textContent = 'Replay in progress...';
      
      setTimeout(() => {{
        statusLabel.textContent = 'Hover to preview, click to pin an event.';
      }}, Math.max(1200, dominoElements.length * 150));
    }});

    window.addEventListener('resize', () => {{
      if (activeDomino) {{
        activeDomino.scrollIntoView({{ behavior: 'smooth', inline: 'center', block: 'nearest' }});
      }}
    }});

    buildDominos();
  </script>
</body>
</html>
"""

timeline_html = f"""
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <style>
    body {{
      margin: 0;
      font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
      background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
      overflow: hidden;
    }}

    #lens {{
      position: fixed;
      top: 80px;
      left: 50%;
      transform: translateX(-50%);
      width: {st.session_state["lens_size"]}px;
      height: {st.session_state["lens_size"]}px;
      border-radius: 50%;
      border: 2px solid transparent;
      background: transparent;
      background-color: transparent !important;
      mix-blend-mode: normal;
      filter: none;
      box-shadow: none;
      pointer-events: none;
      z-index: 1000;
      overflow: visible;
    }}

    #lens::before,
    #lens::after {{
      content: none;
    }}

    #timeline-wrapper {{
      position: relative;
      width: 100%;
      overflow-x: auto;
      overflow-y: hidden;
      white-space: nowrap;
      padding: 180px 50vw 60px;
      box-sizing: border-box;
      scroll-behavior: smooth;
      -ms-overflow-style: none;
      scrollbar-width: none;
    }}

    #timeline-wrapper::-webkit-scrollbar {{
      display: none;
    }}

    .event {{
      display: inline-block;
      width: 120px;
      margin: 0 60px;
      text-align: center;
      color: #ffffff;
      transition: transform 0.6s cubic-bezier(0.68, -0.55, 0.265, 1.55), box-shadow 0.35s ease;
      transform-origin: center center;
      transform: translateY(0) scale(1);
      z-index: 1500;
    }}

    .bubble {{
      width: 70px;
      height: 70px;
      border-radius: 50%;
      margin: 0 auto;
      display: flex;
      align-items: center;
      justify-content: center;
      position: relative;
      overflow: hidden;
      border: 2px solid rgba(255, 255, 255, 0.4);
      background: var(--color);
      box-shadow: 0 4px 15px rgba(0, 0, 0, 0.3), 0 0 20px rgba(102, 126, 234, 0.2);
      transition: border-color 0.35s ease, box-shadow 0.35s ease, transform 0.35s ease;
    }}

    .bubble-title {{
      font-size: 0.55rem;
      font-weight: 600;
      color: #ffffff;
      text-align: center;
      padding: 3px;
      line-height: 1.0;
      text-transform: none;
      letter-spacing: 0.01em;
      overflow: hidden;
      text-overflow: ellipsis;
      display: -webkit-box;
      -webkit-line-clamp: 2;
      -webkit-box-orient: vertical;
      word-break: break-word;
      max-width: 60px;
      max-height: 45px;
    }}

    .bubble::before {{
      content: "";
      position: absolute;
      inset: -40%;
      background: var(--color);
      filter: blur(25px);
      opacity: 0.8;
      transform: scale(1.2);
      transition: opacity 0.35s ease;
      pointer-events: none;
    }}

    .bubble::after {{
      content: "";
      position: absolute;
      inset: 0;
      background: linear-gradient(150deg, rgba(255, 255, 255, 0.9) 0%, rgba(255, 255, 255, 0.3) 55%, rgba(255, 255, 255, 0.1) 100%);
      opacity: 0.95;
      transition: opacity 0.35s ease;
      pointer-events: none;
    }}

    .bubble img {{
      width: 100%;
      height: 100%;
      object-fit: cover;
      filter: blur(8px) saturate(0.55);
      transform: scale(1.12);
      transition: filter 0.35s ease, transform 0.35s ease;
      animation: tremor 1.4s infinite ease-in-out;
    }}

    .label {{
      margin-top: 10px;
      font-weight: bold;
      color: {st.session_state["event_title_color"]};
      transition: color 0.35s ease;
      font-family: {st.session_state["event_title_font"]};
      font-size: {st.session_state["event_title_size"]}px;
      white-space: normal;
      word-wrap: break-word;
      line-height: 1.2;
    }}

    .date {{
      margin-top: 4px;
      font-size: {st.session_state["event_date_size"]}px;
      color: {st.session_state["event_date_color"]};
      letter-spacing: 0.4px;
      transition: color 0.35s ease;
    }}

    .event.focused {{
      transform: scale(2.5);
      z-index: 10;
    }}

    .event.focused .bubble {{
      border-color: rgba(255, 255, 255, 0.95);
      box-shadow: 0 8px 25px rgba(0, 0, 0, 0.4), 0 0 30px rgba(102, 126, 234, 0.4);
    }}

    .event.focused .bubble::before,
    .event.focused .bubble::after {{
      opacity: 0;
    }}

    .event.focused .bubble img {{
      filter: blur(0) saturate(1);
      transform: scale(1);
      animation: none;
    }}

    .event.hovering {{
      transform: translateY(-140px) scale(calc({st.session_state["lens_size"]} / 90));
      z-index: 2000;
    }}

    .event.hovering .bubble {{
      border-color: rgba(255, 255, 255, 0.9);
      box-shadow: 0 8px 25px rgba(0, 0, 0, 0.4), 0 0 25px rgba(102, 126, 234, 0.3);
      transform: scale(1.1);
    }}

    .event.hovering .label {{
      color: #ffffff;
    }}

    .event.hovering .date {{
      color: rgba(255, 255, 255, 0.75);
    }}

    .event.hovering.focused {{
      transform: translateY(-140px) scale(calc({st.session_state["lens_size"]} / 60));
      z-index: 2500;
    }}

    @keyframes tremor {{
      0%, 100% {{ transform: scale(1.12) translate(0px, 0px); }}
      20% {{ transform: scale(1.12) translate(-1.5px, 1.2px); }}
      40% {{ transform: scale(1.12) translate(1.4px, -1px); }}
      60% {{ transform: scale(1.12) translate(-1px, -1.6px); }}
      80% {{ transform: scale(1.12) translate(1.6px, 1px); }}
      100% {{ transform: scale(1.12) translate(1.8px, 1px); }}
    }}
  </style>
</head>
<body>
  <div id="lens"></div>
  <div id="timeline-wrapper"></div>
  <script>
    const events = {events_json};
    events.sort((a, b) => new Date(a.date) - new Date(b.date));
    const wrapper = document.getElementById('timeline-wrapper');
    const eventElements = [];
    const magnificationTimeouts = new Map();

    events.forEach((ev, idx) => {{
      const eventDiv = document.createElement('div');
      eventDiv.className = 'event';
      eventDiv.dataset.id = ev.id;

      const hue = (idx * 137.5) % 360;
      const color = `hsl(${{hue}}, 70%, 50%)`;
      eventDiv.style.setProperty('--color', color);

      const bubbleContent = ev.image 
        ? `<img src=\"data:image/png;base64,${{ev.image}}\" style=\"width:100%;height:100%;object-fit:cover;border-radius:50%;\">`
        : `<span class=\"bubble-title\">${{ev.title || 'Event'}}</span>`;
      const dateLabel = new Date(ev.date).toLocaleDateString(undefined, {{ year: 'numeric', month: 'short', day: 'numeric' }});
      eventDiv.innerHTML = `
        <div class=\"bubble\">${{bubbleContent}}</div>
        <div class=\"label\">${{ev.title}}</div>
        <div class=\"date\">${{dateLabel}}</div>
      `;

      eventDiv.addEventListener('mouseenter', () => {{
        // Clear any existing timeout for this event
        if (magnificationTimeouts.has(eventDiv.dataset.id)) {{
          clearTimeout(magnificationTimeouts.get(eventDiv.dataset.id));
          magnificationTimeouts.delete(eventDiv.dataset.id);
        }}
        
        // Set a timeout to add the hovering class after the specified delay
        const timeoutId = setTimeout(() => {{
          eventDiv.classList.add('hovering');
          const targetScroll = getScrollLeftForEvent(eventDiv);
          wrapper.scrollTo({{ left: targetScroll, behavior: 'smooth' }});
          magnificationTimeouts.delete(eventDiv.dataset.id);
        }}, {st.session_state["lens_duration"]} * 1000);
        
        magnificationTimeouts.set(eventDiv.dataset.id, timeoutId);
      }});

      eventDiv.addEventListener('mouseleave', () => {{
        // Clear any pending timeout and immediately remove hovering
        if (magnificationTimeouts.has(eventDiv.dataset.id)) {{
          clearTimeout(magnificationTimeouts.get(eventDiv.dataset.id));
          magnificationTimeouts.delete(eventDiv.dataset.id);
        }}
        
        eventDiv.classList.remove('hovering');
      }});

      wrapper.appendChild(eventDiv);
      eventElements.push(eventDiv);
    }});

    let lensX = window.innerWidth / 2;

    function getScrollLeftForEvent(el) {{
      const wrapperRect = wrapper.getBoundingClientRect();
      const targetCenter = el.offsetLeft + el.offsetWidth / 2;
      const lensOffset = lensX - wrapperRect.left;
      return Math.max(0, targetCenter - lensOffset);
    }}

    function updateFocus() {{
      document.querySelectorAll('.event').forEach(el => {{
        const rect = el.getBoundingClientRect();
        const center = rect.left + rect.width / 2;
        const lensRadius = {st.session_state["lens_size"]} / 2;
        if (Math.abs(center - lensX) < lensRadius * 0.8) {{
          el.classList.add('focused');
        }} else {{
          el.classList.remove('focused');
        }}
      }});
    }}

    updateFocus();

    wrapper.addEventListener('scroll', () => {{
      updateFocus();
    }});

    window.addEventListener('resize', () => {{
      lensX = window.innerWidth / 2;
      updateFocus();
    }});

  </script>
</body>
</html>
"""

if view_mode == "Domino":
    st.components.v1.html(domino_html, height=620, scrolling=False)
else:
    st.components.v1.html(timeline_html, height=400, scrolling=True)
